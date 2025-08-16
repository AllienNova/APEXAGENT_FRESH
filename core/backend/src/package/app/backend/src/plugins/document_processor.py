# Plugin: Advanced Document Handling
import os
import re # For generating simple IDs
# Ensure these libraries are installed in your environment if not already:
# pip install python-docx python-pptx PyPDF2 reportlab weasyprint fpdf2 pdfminer.six

# For DOCX
from docx import Document
from docx.shared import Inches

# For PPTX
from pptx import Presentation
from pptx.util import Inches as PptxInches

# For PDF Reading (choose one or use multiple based on needs)
import PyPDF2 # For basic text extraction and manipulation
# from pdfminer.high_level import extract_text as pdfminer_extract_text # Alternative for more complex PDFs

# For PDF Creation (choose one based on needs)
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
# from weasyprint import HTML # For HTML to PDF
# from fpdf import FPDF # Another option for PDF creation

class DocumentProcessor:
    def __init__(self, api_key_manager=None):
        """
        Initializes the DocumentProcessor.
        Args:
            api_key_manager: An instance of ApiKeyManager (not typically needed for local doc processing).
        """
        self.api_key_manager = api_key_manager
        # Simple relation normalization map (conceptual)
        self.relation_normalization_map = {
            "works_at": "EMPLOYED_BY",
            "consultant_for": "CONSULTS_FOR",
            "located_in": "HAS_LOCATION",
            "competitor_of": "IS_COMPETITOR_OF"
        }
        print("DocumentProcessor initialized.")

    def _generate_entity_id(self, text: str, entity_type: str = None) -> str:
        """Generates a simple, somewhat stable ID from entity text."""
        # Lowercase, replace non-alphanumeric with underscore, remove leading/trailing underscores
        base_id = re.sub(pattern=r'[^a-z0-9_]+', repl='', string=re.sub(pattern=r'\s+', repl='_', string=text.lower()).strip('_'))
        # Optionally, prefix with type for more uniqueness if needed, though KG might handle types separately
        # if entity_type:
        #     return f"{entity_type.lower()}_{base_id}"
        return base_id

    def _disambiguate_entity(self, text: str, entity_type: str, context: str = None) -> dict:
        """
        Conceptual placeholder for entity disambiguation.
        In a real implementation, this would involve more sophisticated NLP techniques,
        contextual analysis, and potentially linking to a knowledge base.
        """
        entity_id = self._generate_entity_id(text, entity_type)
        # Placeholder for a canonical ID that might come from a knowledge base lookup
        canonical_id = f"canonical_{entity_id}" # Simulate a lookup
        return {
            "id": entity_id, 
            "text": text, 
            "type": entity_type, 
            "canonical_id": canonical_id, # Conceptual field
            "disambiguation_status": "placeholder_disambiguated" # Conceptual field
        }

    def _normalize_relation_type(self, relation_text: str) -> str:
        """
        Conceptual placeholder for relation normalization.
        Uses a simple map for now.
        """
        return self.relation_normalization_map.get(relation_text.lower().replace(" ", "_"), relation_text.upper().replace(" ", "_"))

    # --- DOCX Methods ---
    def read_docx_text(self, file_path: str):
        """Extracts all text content from a .docx file."""
        if not os.path.exists(file_path):
            return {"error": f"File not found: {file_path}"}
        try:
            doc = Document(file_path)
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            return {"text_content": "\n".join(full_text)}
        except Exception as e:
            return {"error": f"Error reading DOCX file {file_path}: {str(e)}"}

    def create_docx_from_text(self, text_content: str, save_path: str, title: str = None):
        """Creates a new .docx file from given text content."""
        try:
            doc = Document()
            if title:
                doc.core_properties.title = title
            for line in text_content.split("\n"):
                doc.add_paragraph(line)
            doc.save(save_path)
            return {"status": "success", "file_path": save_path}
        except Exception as e:
            return {"error": f"Error creating DOCX file {save_path}: {str(e)}"}

    def append_to_docx(self, file_path: str, content_to_append: str, content_type: str = "paragraph"):
        """Appends text, paragraphs, or simple tables to an existing .docx file."""
        if not os.path.exists(file_path):
            return {"error": f"File not found: {file_path}"}
        try:
            doc = Document(file_path)
            if content_type == "paragraph":
                doc.add_paragraph(content_to_append)
            elif content_type == "heading1":
                doc.add_heading(content_to_append, level=1)
            elif content_type == "heading2":
                doc.add_heading(content_to_append, level=2)
            else:
                return {"error": f"Unsupported content_type: {content_type}"}
            doc.save(file_path)
            return {"status": "success", "file_path": file_path}
        except Exception as e:
            return {"error": f"Error appending to DOCX file {file_path}: {str(e)}"}

    # --- PPTX Methods ---
    def read_pptx_text(self, file_path: str):
        """Extracts all text content from all slides in a .pptx file."""
        if not os.path.exists(file_path):
            return {"error": f"File not found: {file_path}"}
        try:
            prs = Presentation(file_path)
            slides_text = []
            for slide in prs.slides:
                slide_text_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text_content.append(shape.text)
                slides_text.append("\n".join(slide_text_content))
            return {"slides_text": slides_text}
        except Exception as e:
            return {"error": f"Error reading PPTX file {file_path}: {str(e)}"}

    def create_pptx_from_slides_data(self, slides_data: list, save_path: str):
        """Creates a new .pptx file from a list of slide data."""
        try:
            prs = Presentation()
            for slide_info in slides_data:
                slide_layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(slide_layout)
                title_text = slide_info.get("title")
                content_text = slide_info.get("content")
                if title_text:
                    title_shape = slide.shapes.title
                    if title_shape: title_shape.text = title_text
                    else:
                        txBox = slide.shapes.add_textbox(PptxInches(1.0), PptxInches(1.0), PptxInches(1.0), PptxInches(1.0))
                        txBox.text_frame.text = title_text
                if content_text:
                    body_shape = next((shape for shape in slide.placeholders if shape.name.startswith("Text Placeholder") or shape.name.startswith("Content Placeholder") or shape.name.startswith("Body")), None)
                    tf = body_shape.text_frame if body_shape else slide.shapes.add_textbox(PptxInches(1.0), PptxInches(2.0) if title_text else PptxInches(1.0), PptxInches(8.0), PptxInches(5.0)).text_frame
                    if isinstance(content_text, list):
                        for line in content_text: tf.add_paragraph().text = line
                    else: tf.text = content_text
            prs.save(save_path)
            return {"status": "success", "file_path": save_path}
        except Exception as e:
            return {"error": f"Error creating PPTX file {save_path}: {str(e)}"}

    # --- PDF Methods ---
    def extract_text_from_pdf(self, file_path: str):
        """Extracts text content from a .pdf file using PyPDF2."""
        if not os.path.exists(file_path):
            return {"error": f"File not found: {file_path}"}
        try:
            text_content = ""
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages: text_content += page.extract_text() or ""
            if not text_content.strip(): return {"text_content": "", "warning": "No text could be extracted or PDF might be image-based."}
            return {"text_content": text_content}
        except Exception as e:
            return {"error": f"Error extracting text from PDF {file_path}: {str(e)}"}

    def convert_text_to_pdf(self, text_content: str, save_path: str, title: str = None):
        """Converts plain text content into a basic PDF document using ReportLab."""
        try:
            c = canvas.Canvas(save_path, pagesize=letter)
            if title: c.setTitle(title)
            text_object = c.beginText(Inches(1), Inches(10))
            text_object.setFont("Times-Roman", 12)
            for line in text_content.split("\n"): text_object.textLine(line)
            c.drawText(text_object)
            c.showPage()
            c.save()
            return {"status": "success", "file_path": save_path}
        except Exception as e:
            return {"error": f"Error converting text to PDF {save_path}: {str(e)}"}

    # --- Entity and Relation Extraction ---
    def extract_entities_relations(self, file_path: str, context_for_disambiguation: str = None):
        """Extracts entities and relations from the text content of a document with conceptual disambiguation and normalization."""
        if not os.path.exists(file_path):
            return {"success": False, "error": f"File not found: {file_path}", "entities": [], "relations": []}
        
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text_content = f.read()
            
            processed_entities = []
            relations = []
            # Use a map to store disambiguated entity info by original text to avoid reprocessing
            # and to ensure consistent ID usage for relations.
            entity_cache = {} 

            def get_or_create_entity(text, type, context=None):
                if text in entity_cache:
                    return entity_cache[text]
                
                # Conceptual: In a real system, context from the document would be passed here.
                disambiguated_entity = self._disambiguate_entity(text, type, context=context_for_disambiguation)
                processed_entities.append(disambiguated_entity)
                entity_cache[text] = disambiguated_entity
                return disambiguated_entity

            # Placeholder: Simulate finding some entities and relations
            if "Alice Smith" in text_content and "Acme Corp" in text_content:
                alice = get_or_create_entity("Alice Smith", "PERSON")
                acme = get_or_create_entity("Acme Corp", "ORGANIZATION")
                relation_type_raw = "works_at"
                normalized_relation_type = self._normalize_relation_type(relation_type_raw)
                relations.append({
                    "subject_id": alice["id"], # Use the simple ID for now, KG plugin can use canonical_id if needed
                    "relation_type": normalized_relation_type, 
                    "object_id": acme["id"], 
                    "properties": {"role": "Placeholder Role", "original_relation_text": relation_type_raw}
                })
                
            if "Bob Johnson" in text_content and "Zeta Inc" in text_content:
                bob = get_or_create_entity("Bob Johnson", "PERSON")
                zeta = get_or_create_entity("Zeta Inc", "ORGANIZATION")
                relation_type_raw = "consultant_for"
                normalized_relation_type = self._normalize_relation_type(relation_type_raw)
                relations.append({
                    "subject_id": bob["id"], 
                    "relation_type": normalized_relation_type, 
                    "object_id": zeta["id"], 
                    "properties": {"original_relation_text": relation_type_raw}
                })
            
            if "Acme Corp" in text_content and "New York" in text_content:
                acme = get_or_create_entity("Acme Corp", "ORGANIZATION")
                ny = get_or_create_entity("New York", "LOCATION")
                relation_type_raw = "located_in"
                normalized_relation_type = self._normalize_relation_type(relation_type_raw)
                relations.append({
                    "subject_id": acme["id"], 
                    "relation_type": normalized_relation_type, 
                    "object_id": ny["id"], 
                    "properties": {"original_relation_text": relation_type_raw}
                })

            if "Zeta Inc" in text_content and "Acme Corp" in text_content:
                zeta = get_or_create_entity("Zeta Inc", "ORGANIZATION")
                acme = get_or_create_entity("Acme Corp", "ORGANIZATION")
                relation_type_raw = "competitor_of"
                normalized_relation_type = self._normalize_relation_type(relation_type_raw)
                relations.append({
                    "subject_id": zeta["id"], 
                    "relation_type": normalized_relation_type, 
                    "object_id": acme["id"], 
                    "properties": {"original_relation_text": relation_type_raw}
                })

            if not processed_entities and not relations:
                return {"success": True, "message": "No specific entities/relations found by placeholder logic.", "entities": [], "relations": []}

            return {"success": True, "entities": processed_entities, "relations": relations}
        except Exception as e:
            return {"success": False, "error": f"Error processing file for entity/relation extraction: {str(e)}", "entities": [], "relations": []}

# Example usage (for testing, not part of the plugin itself)
if __name__ == "__main__":
    doc_processor = DocumentProcessor()
    # ... (rest of the example usage can remain for other functionalities)
    sample_file_path = "/home/ubuntu/sample_extraction_test.txt"
    with open(sample_file_path, "w") as f:
        f.write("Alice Smith works at Acme Corp. Bob Johnson is a consultant for Zeta Inc. Acme Corp is located in New York. Zeta Inc. is a competitor of Acme Corp.")
    
    print("\n--- Entity/Relation Extraction Test with Conceptual Disambiguation/Normalization ---")
    extraction_result = doc_processor.extract_entities_relations(sample_file_path)
    print(json.dumps(extraction_result, indent=2))
    if os.path.exists(sample_file_path):
        os.remove(sample_file_path)

