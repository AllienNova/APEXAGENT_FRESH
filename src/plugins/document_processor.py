# Plugin: Advanced Document Handling
import os
import re # For generating simple IDs
# Ensure these libraries are installed in your environment if not already:
# pip install python-docx python-pptx PyPDF2 reportlab weasyprint fpdf2 pdfminer.six

# For DOCX
from docx import Document
from docx.shared import Inches

# For PPTX
from pptx import Presentation
from pptx.util import Inches as PptxInches

# For PDF Reading (choose one or use multiple based on needs)
import PyPDF2 # For basic text extraction and manipulation
# from pdfminer.high_level import extract_text as pdfminer_extract_text # Alternative for more complex PDFs

# For PDF Creation (choose one based on needs)
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
# from weasyprint import HTML # For HTML to PDF
# from fpdf import FPDF # Another option for PDF creation

class DocumentProcessor:
    def __init__(self, api_key_manager=None):
        """
        Initializes the DocumentProcessor.
        Args:
            api_key_manager: An instance of ApiKeyManager (not typically needed for local doc processing).
        """
        self.api_key_manager = api_key_manager
        # Simple relation normalization map (conceptual)
        self.relation_normalization_map = {
            "works_at": "EMPLOYED_BY",
            "consultant_for": "CONSULTS_FOR",
            "located_in": "HAS_LOCATION",
            "competitor_of": "IS_COMPETITOR_OF"
        }
        print("DocumentProcessor initialized.")

    def _generate_entity_id(self, text: str, entity_type: str = None) -> str:
        """Generates a simple, somewhat stable ID from entity text."""
        # Lowercase, replace non-alphanumeric with underscore, remove leading/trailing underscores
        base_id = re.sub(pattern=r'[^a-z0-9_]+', repl='', string=re.sub(pattern=r'\s+', repl='_', string=text.lower()).strip('_'))
        # Optionally, prefix with type for more uniqueness if needed, though KG might handle types separately
        # if entity_type:
        #     return f"{entity_type.lower()}_{base_id}"
        return base_id

    def _disambiguate_entity(self, text: str, entity_type: str, context: str = None) -> dict:
        """
        Conceptual placeholder for entity disambiguation.
        In a real implementation, this would involve more sophisticated NLP techniques,
        contextual analysis, and potentially linking to a knowledge base.
        """
        entity_id = self._generate_entity_id(text, entity_type)
        # Placeholder for a canonical ID that might come from a knowledge base lookup
        canonical_id = f"canonical_{entity_id}" # Simulate a lookup
        return {
            "id": entity_id, 
            "text": text, 
            "type": entity_type, 
            "canonical_id": canonical_id, # Conceptual field
            "disambiguation_status": "placeholder_disambiguated" # Conceptual field
        }

    def _normalize_relation_type(self, relation_text: str) -> str:
        """
        Conceptual placeholder for relation normalization.
        Uses a simple map for now.
        """
        return self.relation_normalization_map.get(relation_text.lower().replace(" ", "_"), relation_text.upper().replace(" ", "_"))

    # --- DOCX Methods ---
    def read_docx_text(self, file_path: str):
        """Extracts all text content from a .docx file."""
        if not os.path.exists(file_path):
            return {"error": f"File not found: {file_path}"}
        try:
            doc = Document(file_path)
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            return {"text_content": "\n".join(full_text)}
        except Exception as e:
            return {"error": f"Error reading DOCX file {file_path}: {str(e)}"}

    def create_docx_from_text(self, text_content: str, save_path: str, title: str = None):
        """Creates a new .docx file from given text content."""
        try:
            doc = Document()
            if title:
                doc.core_properties.title = title
            for line in text_content.split("\n"):
                doc.add_paragraph(line)
            doc.save(save_path)
            return {"status": "success", "file_path": save_path}
        except Exception as e:
            return {"error": f"Error creating DOCX file {save_path}: {str(e)}"}

    def append_to_docx(self, file_path: str, content_to_append: str, content_type: str = "paragraph"):
        """Appends text, paragraphs, or simple tables to an existing .docx file."""
        if not os.path.exists(file_path):
            return {"error": f"File not found: {file_path}"}
        try:
            doc = Document(file_path)
            if content_type == "paragraph":
                doc.add_paragraph(content_to_append)
            elif content_type == "heading1":
                doc.add_heading(content_to_append, level=1)
            elif content_type == "heading2":
                doc.add_heading(content_to_append, level=2)
            else:
                return {"error": f"Unsupported content_type: {content_type}"}
            doc.save(file_path)
            return {"status": "success", "file_path": file_path}
        except Exception as e:
            return {"error": f"Error appending to DOCX file {file_path}: {str(e)}"}

    # --- PPTX Methods ---
    def read_pptx_text(self, file_path: str):
        """Extracts all text content from all slides in a .pptx file."""
        if not os.path.exists(file_path):
            return {"error": f"File not found: {file_path}"}
        try:
            prs = Presentation(file_path)
            slides_text = []
            for slide in prs.slides:
                slide_text_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text_content.append(shape.text)
                slides_text.append("\n".join(slide_text_content))
            return {"slides_text": slides_text}
        except Exception as e:
            return {"error": f"Error reading PPTX file {file_path}: {str(e)}"}

    def create_pptx_from_slides_data(self, slides_data: list, save_path: str):
        """Creates a new .pptx file from a list of slide data."""
        try:
            prs = Presentation()
            for slide_info in slides_data:
                slide_layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(slide_layout)
                title_text = slide_info.get("title")
                content_text = slide_info.get("content")
                if title_text:
                    title_shape = slide.shapes.title
                    if title_shape: title_shape.text = title_text
                    else:
                        txBox = slide.shapes.add_textbox(PptxInches(1.0), PptxInches(1.0), PptxInches(1.0), PptxInches(1.0))
                        txBox.text_frame.text = title_text
                if content_text:
                    body_shape = next((shape for shape in slide.placeholders if shape.name.startswith("Text Placeholder") or shape.name.startswith("Content Placeholder") or shape.name.startswith("Body")), None)
                    tf = body_shape.text_frame if body_shape else slide.shapes.add_textbox(PptxInches(1.0), PptxInches(2.0) if title_text else PptxInches(1.0), PptxInches(8.0), PptxInches(5.0)).text_frame
                    if isinstance(content_text, list):
                        for line in content_text: tf.add_paragraph().text = line
                    else: tf.text = content_text
            prs.save(save_path)
            return {"status": "success", "file_path": save_path}
        except Exception as e:
            return {"error": f"Error creating PPTX file {save_path}: {str(e)}"}

    # --- PDF Methods ---
    def extract_text_from_pdf(self, file_path: str):
        """Extracts text content from a .pdf file using PyPDF2."""
        if not os.path.exists(file_path):
            return {"error": f"File not found: {file_path}"}
        try:
            text_content = ""
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages: text_content += page.extract_text() or ""
            if not text_content.strip(): return {"text_content": "", "warning": "No text could be extracted or PDF might be image-based."}
            return {"text_content": text_content}
        except Exception as e:
            return {"error": f"Error extracting text from PDF {file_path}: {str(e)}"}

    def convert_text_to_pdf(self, text_content: str, save_path: str, title: str = None):
        """Converts plain text content into a basic PDF document using ReportLab."""
        try:
            c = canvas.Canvas(save_path, pagesize=letter)
            if title: c.setTitle(title)
            text_object = c.beginText(Inches(1), Inches(10))
            text_object.setFont("Times-Roman", 12)
            for line in text_content.split("\n"): text_object.textLine(line)
            c.drawText(text_object)
            c.showPage()
            c.save()
            return {"status": "success", "file_path": save_path}
        except Exception as e:
            return {"error": f"Error converting text to PDF {save_path}: {str(e)}"}

    # --- Entity and Relation Extraction ---
    def extract_entities_relations(self, file_path: str, context_for_disambiguation: str = None):
        """Extracts entities and relations from the text content of a document with disambiguation and normalization."""
        if not os.path.exists(file_path):
            return {"success": False, "error": f"File not found: {file_path}", "entities": [], "relations": []}
        
        try:
            # Determine file type and extract text accordingly
            file_ext = os.path.splitext(file_path)[1].lower()
            text_content = ""
            
            if file_ext == ".txt":
                with open(file_path, "r", encoding="utf-8") as f:
                    text_content = f.read()
            elif file_ext == ".docx":
                result = self.read_docx_text(file_path)
                if "error" in result:
                    return {"success": False, "error": result["error"], "entities": [], "relations": []}
                text_content = result.get("text_content", "")
            elif file_ext == ".pdf":
                result = self.extract_text_from_pdf(file_path)
                if "error" in result:
                    return {"success": False, "error": result["error"], "entities": [], "relations": []}
                text_content = result.get("text_content", "")
            elif file_ext == ".pptx":
                result = self.read_pptx_text(file_path)
                if "error" in result:
                    return {"success": False, "error": result["error"], "entities": [], "relations": []}
                text_content = "\n".join(result.get("slides_text", []))
            else:
                return {"success": False, "error": f"Unsupported file type: {file_ext}", "entities": [], "relations": []}
            
            # Initialize NLP components for entity and relation extraction
            try:
                import spacy
                from spacy.matcher import DependencyMatcher
                
                # Load appropriate spaCy model
                try:
                    nlp = spacy.load("en_core_web_lg")  # Larger model for better entity recognition
                except OSError:
                    # Fall back to smaller model if large one isn't available
                    try:
                        nlp = spacy.load("en_core_web_sm")
                    except OSError:
                        # If no model is installed, install it automatically
                        import subprocess
                        try:
                            subprocess.check_call([sys.executable, "-m", "spacy", "download", "en_core_web_sm"])
                            nlp = spacy.load("en_core_web_sm")
                        except Exception as e:
                            return {
                                "success": False, 
                                "error": f"Failed to install required spaCy models: {str(e)}", 
                                "entities": [], 
                                "relations": []
                            }
            except ImportError:
                # If spaCy isn't installed, install it automatically
                import subprocess
                import sys
                try:
                    subprocess.check_call([sys.executable, "-m", "pip", "install", "spacy"])
                    subprocess.check_call([sys.executable, "-m", "spacy", "download", "en_core_web_sm"])
                    import spacy
                    from spacy.matcher import DependencyMatcher
                    nlp = spacy.load("en_core_web_sm")
                except Exception as e:
                    # Fall back to regex-based extraction if installation fails
                    import re
                    nlp = None
            
            processed_entities = []
            relations = []
            # Use a map to store disambiguated entity info by original text to avoid reprocessing
            # and to ensure consistent ID usage for relations.
            entity_cache = {} 

            def get_or_create_entity(text, type, context=None):
                if text in entity_cache:
                    return entity_cache[text]
                
                disambiguated_entity = self._disambiguate_entity(text, type, context=context_for_disambiguation)
                processed_entities.append(disambiguated_entity)
                entity_cache[text] = disambiguated_entity
                return disambiguated_entity

            # If spaCy is available, use it for sophisticated entity and relation extraction
            if nlp:
                # Process the text with spaCy
                doc = nlp(text_content)
                
                # Extract named entities
                for ent in doc.ents:
                    # Map spaCy entity types to our system's types
                    entity_type_map = {
                        "PERSON": "PERSON",
                        "ORG": "ORGANIZATION",
                        "GPE": "LOCATION",
                        "LOC": "LOCATION",
                        "PRODUCT": "PRODUCT",
                        "DATE": "DATE",
                        "TIME": "TIME",
                        "MONEY": "MONEY",
                        "PERCENT": "PERCENT",
                        "FACILITY": "FACILITY",
                        "NORP": "GROUP",  # Nationalities, religious or political groups
                        "EVENT": "EVENT",
                        "WORK_OF_ART": "CREATIVE_WORK",
                        "LAW": "LAW",
                        "LANGUAGE": "LANGUAGE"
                    }
                    
                    entity_type = entity_type_map.get(ent.label_, "MISC")
                    get_or_create_entity(ent.text, entity_type)
                
                # Extract relations using dependency parsing
                # Create a dependency matcher for common relation patterns
                matcher = DependencyMatcher(nlp.vocab)
                
                # Pattern for "works at" relation
                works_at_pattern = [
                    {
                        "RIGHT_ID": "person",
                        "RIGHT_ATTRS": {"ENT_TYPE": "PERSON"}
                    },
                    {
                        "LEFT_ID": "person",
                        "REL_OP": ">",
                        "RIGHT_ID": "verb",
                        "RIGHT_ATTRS": {"LEMMA": {"IN": ["work", "employ"]}}
                    },
                    {
                        "LEFT_ID": "verb",
                        "REL_OP": ">",
                        "RIGHT_ID": "prep",
                        "RIGHT_ATTRS": {"DEP": "prep", "LEMMA": {"IN": ["at", "for", "by", "with"]}}
                    },
                    {
                        "LEFT_ID": "prep",
                        "REL_OP": ">",
                        "RIGHT_ID": "org",
                        "RIGHT_ATTRS": {"ENT_TYPE": "ORG"}
                    }
                ]
                
                # Pattern for "located in" relation
                located_in_pattern = [
                    {
                        "RIGHT_ID": "entity",
                        "RIGHT_ATTRS": {"ENT_TYPE": {"IN": ["ORG", "FACILITY", "PRODUCT"]}}
                    },
                    {
                        "LEFT_ID": "entity",
                        "REL_OP": ">",
                        "RIGHT_ID": "verb",
                        "RIGHT_ATTRS": {"LEMMA": {"IN": ["locate", "base", "situate", "headquarter"]}}
                    },
                    {
                        "LEFT_ID": "verb",
                        "REL_OP": ">",
                        "RIGHT_ID": "prep",
                        "RIGHT_ATTRS": {"DEP": "prep", "LEMMA": "in"}
                    },
                    {
                        "LEFT_ID": "prep",
                        "REL_OP": ">",
                        "RIGHT_ID": "location",
                        "RIGHT_ATTRS": {"ENT_TYPE": {"IN": ["GPE", "LOC"]}}
                    }
                ]
                
                # Add patterns to matcher
                matcher.add("WORKS_AT", [works_at_pattern])
                matcher.add("LOCATED_IN", [located_in_pattern])
                
                # Find matches
                matches = matcher(doc)
                
                # Process matches to extract relations
                for match_id, token_ids in matches:
                    pattern_name = nlp.vocab.strings[match_id]
                    
                    if pattern_name == "WORKS_AT":
                        # Extract the person and organization
                        person_token = doc[token_ids[0]]
                        org_token = doc[token_ids[3]]
                        
                        # Get or create entities
                        person = get_or_create_entity(person_token.text, "PERSON")
                        org = get_or_create_entity(org_token.text, "ORGANIZATION")
                        
                        # Create relation
                        relation_type_raw = "works_at"
                        normalized_relation_type = self._normalize_relation_type(relation_type_raw)
                        relations.append({
                            "subject_id": person["id"],
                            "relation_type": normalized_relation_type,
                            "object_id": org["id"],
                            "properties": {
                                "confidence": 0.85,
                                "extraction_method": "dependency_pattern",
                                "original_relation_text": relation_type_raw
                            }
                        })
                    
                    elif pattern_name == "LOCATED_IN":
                        # Extract the entity and location
                        entity_token = doc[token_ids[0]]
                        location_token = doc[token_ids[3]]
                        
                        # Get or create entities
                        entity_type = "ORGANIZATION" if entity_token.ent_type_ == "ORG" else "FACILITY"
                        entity = get_or_create_entity(entity_token.text, entity_type)
                        location = get_or_create_entity(location_token.text, "LOCATION")
                        
                        # Create relation
                        relation_type_raw = "located_in"
                        normalized_relation_type = self._normalize_relation_type(relation_type_raw)
                        relations.append({
                            "subject_id": entity["id"],
                            "relation_type": normalized_relation_type,
                            "object_id": location["id"],
                            "properties": {
                                "confidence": 0.9,
                                "extraction_method": "dependency_pattern",
                                "original_relation_text": relation_type_raw
                            }
                        })
                
                # Additional relation extraction using co-occurrence and heuristics
                sentences = list(doc.sents)
                
                for sentence in sentences:
                    sentence_entities = [ent for ent in doc.ents if ent.start >= sentence.start and ent.end <= sentence.end]
                    
                    # Look for competitor relationships between organizations
                    orgs = [ent for ent in sentence_entities if ent.label_ == "ORG"]
                    if len(orgs) >= 2:
                        # Check for competitor keywords
                        competitor_keywords = ["competitor", "compete", "rival", "competition"]
                        sentence_text = sentence.text.lower()
                        
                        if any(keyword in sentence_text for keyword in competitor_keywords):
                            for i in range(len(orgs)):
                                for j in range(i+1, len(orgs)):
                                    org1 = get_or_create_entity(orgs[i].text, "ORGANIZATION")
                                    org2 = get_or_create_entity(orgs[j].text, "ORGANIZATION")
                                    
                                    relation_type_raw = "competitor_of"
                                    normalized_relation_type = self._normalize_relation_type(relation_type_raw)
                                    
                                    # Create bidirectional competitor relationship
                                    relations.append({
                                        "subject_id": org1["id"],
                                        "relation_type": normalized_relation_type,
                                        "object_id": org2["id"],
                                        "properties": {
                                            "confidence": 0.75,
                                            "extraction_method": "co_occurrence",
                                            "original_relation_text": relation_type_raw
                                        }
                                    })
                    
                    # Look for consultant relationships
                    persons = [ent for ent in sentence_entities if ent.label_ == "PERSON"]
                    if len(persons) >= 1 and len(orgs) >= 1:
                        # Check for consultant keywords
                        consultant_keywords = ["consultant", "consult", "advise", "advisor"]
                        sentence_text = sentence.text.lower()
                        
                        if any(keyword in sentence_text for keyword in consultant_keywords):
                            for person in persons:
                                for org in orgs:
                                    person_entity = get_or_create_entity(person.text, "PERSON")
                                    org_entity = get_or_create_entity(org.text, "ORGANIZATION")
                                    
                                    relation_type_raw = "consultant_for"
                                    normalized_relation_type = self._normalize_relation_type(relation_type_raw)
                                    
                                    relations.append({
                                        "subject_id": person_entity["id"],
                                        "relation_type": normalized_relation_type,
                                        "object_id": org_entity["id"],
                                        "properties": {
                                            "confidence": 0.7,
                                            "extraction_method": "keyword_proximity",
                                            "original_relation_text": relation_type_raw
                                        }
                                    })
            
            else:
                # Fallback to regex-based extraction if spaCy is not available
                # Define regex patterns for common entity types
                patterns = {
                    "PERSON": r'\b[A-Z][a-z]+ (?:[A-Z][a-z]+\s?)+',  # Simple pattern for names
                    "ORGANIZATION": r'\b[A-Z][a-z]+ (?:Inc|Corp|LLC|Ltd|Company|Co)\b',  # Simple pattern for organizations
                    "LOCATION": r'\b(?:New York|Los Angeles|Chicago|Houston|Phoenix|Philadelphia|San Antonio|San Diego|Dallas|San Jose|Austin|Jacksonville|Fort Worth|Columbus|San Francisco|Charlotte|Indianapolis|Seattle|Denver|Washington|Boston|El Paso|Nashville|Detroit|Oklahoma City|Portland|Las Vegas|Memphis|Louisville|Baltimore|Milwaukee|Albuquerque|Tucson|Fresno|Sacramento|Mesa|Kansas City|Atlanta|Long Beach|Colorado Springs|Raleigh|Miami|Omaha|Minneapolis|Tulsa|Cleveland|Wichita|Arlington|New Orleans|Bakersfield|Tampa|Honolulu|Aurora|Anaheim|Santa Ana|St\. Louis|Riverside|Corpus Christi|Lexington|Pittsburgh|Anchorage|Stockton|Cincinnati|St\. Paul|Toledo|Greensboro|Newark|Plano|Henderson|Lincoln|Buffalo|Jersey City|Chula Vista|Fort Wayne|Orlando|St\. Petersburg|Chandler|Laredo|Norfolk|Durham|Madison|Lubbock|Irvine|Winston-Salem|Glendale|Garland|Hialeah|Reno|Chesapeake|Gilbert|Baton Rouge|Irving|Scottsdale|North Las Vegas|Fremont|Boise City|Richmond|San Bernardino)\b'  # Common US cities
                }
                
                # Extract entities using regex
                for entity_type, pattern in patterns.items():
                    matches = re.finditer(pattern, text_content)
                    for match in matches:
                        entity_text = match.group(0)
                        get_or_create_entity(entity_text, entity_type)
                
                # Simple relation extraction based on proximity and keywords
                # This is a simplified approach and would be much more sophisticated in a real implementation
                
                # Extract "works at" relations
                works_at_pattern = r'(\b[A-Z][a-z]+ (?:[A-Z][a-z]+\s?)+)\s+(?:works|is employed|is working)\s+(?:at|for|by)\s+(\b[A-Z][a-z]+ (?:Inc|Corp|LLC|Ltd|Company|Co)\b)'
                for match in re.finditer(works_at_pattern, text_content):
                    person_text = match.group(1)
                    org_text = match.group(2)
                    
                    person = get_or_create_entity(person_text, "PERSON")
                    org = get_or_create_entity(org_text, "ORGANIZATION")
                    
                    relation_type_raw = "works_at"
                    normalized_relation_type = self._normalize_relation_type(relation_type_raw)
                    relations.append({
                        "subject_id": person["id"],
                        "relation_type": normalized_relation_type,
                        "object_id": org["id"],
                        "properties": {
                            "confidence": 0.7,
                            "extraction_method": "regex",
                            "original_relation_text": relation_type_raw
                        }
                    })
                
                # Extract "located in" relations
                located_in_pattern = r'(\b[A-Z][a-z]+ (?:Inc|Corp|LLC|Ltd|Company|Co)\b)\s+(?:is located|is based|is headquartered|is situated)\s+in\s+(\b(?:New York|Los Angeles|Chicago|Houston|Phoenix|Philadelphia|San Antonio|San Diego|Dallas|San Jose|Austin|Jacksonville|Fort Worth|Columbus|San Francisco|Charlotte|Indianapolis|Seattle|Denver|Washington|Boston|El Paso|Nashville|Detroit|Oklahoma City|Portland|Las Vegas|Memphis|Louisville|Baltimore|Milwaukee|Albuquerque|Tucson|Fresno|Sacramento|Mesa|Kansas City|Atlanta|Long Beach|Colorado Springs|Raleigh|Miami|Omaha|Minneapolis|Tulsa|Cleveland|Wichita|Arlington|New Orleans|Bakersfield|Tampa|Honolulu|Aurora|Anaheim|Santa Ana|St\. Louis|Riverside|Corpus Christi|Lexington|Pittsburgh|Anchorage|Stockton|Cincinnati|St\. Paul|Toledo|Greensboro|Newark|Plano|Henderson|Lincoln|Buffalo|Jersey City|Chula Vista|Fort Wayne|Orlando|St\. Petersburg|Chandler|Laredo|Norfolk|Durham|Madison|Lubbock|Irvine|Winston-Salem|Glendale|Garland|Hialeah|Reno|Chesapeake|Gilbert|Baton Rouge|Irving|Scottsdale|North Las Vegas|Fremont|Boise City|Richmond|San Bernardino)\b)'
                for match in re.finditer(located_in_pattern, text_content):
                    org_text = match.group(1)
                    location_text = match.group(2)
                    
                    org = get_or_create_entity(org_text, "ORGANIZATION")
                    location = get_or_create_entity(location_text, "LOCATION")
                    
                    relation_type_raw = "located_in"
                    normalized_relation_type = self._normalize_relation_type(relation_type_raw)
                    relations.append({
                        "subject_id": org["id"],
                        "relation_type": normalized_relation_type,
                        "object_id": location["id"],
                        "properties": {
                            "confidence": 0.7,
                            "extraction_method": "regex",
                            "original_relation_text": relation_type_raw
                        }
                    })
                
                # Extract "competitor of" relations
                competitor_pattern = r'(\b[A-Z][a-z]+ (?:Inc|Corp|LLC|Ltd|Company|Co)\b)\s+(?:is a competitor|competes|is a rival|rivals)\s+(?:of|with)\s+(\b[A-Z][a-z]+ (?:Inc|Corp|LLC|Ltd|Company|Co)\b)'
                for match in re.finditer(competitor_pattern, text_content):
                    org1_text = match.group(1)
                    org2_text = match.group(2)
                    
                    org1 = get_or_create_entity(org1_text, "ORGANIZATION")
                    org2 = get_or_create_entity(org2_text, "ORGANIZATION")
                    
                    relation_type_raw = "competitor_of"
                    normalized_relation_type = self._normalize_relation_type(relation_type_raw)
                    relations.append({
                        "subject_id": org1["id"],
                        "relation_type": normalized_relation_type,
                        "object_id": org2["id"],
                        "properties": {
                            "confidence": 0.7,
                            "extraction_method": "regex",
                            "original_relation_text": relation_type_raw
                        }
                    })
                
                # Extract "consultant for" relations
                consultant_pattern = r'(\b[A-Z][a-z]+ (?:[A-Z][a-z]+\s?)+)\s+(?:is a consultant|consults|is an advisor|advises)\s+(?:for|to)\s+(\b[A-Z][a-z]+ (?:Inc|Corp|LLC|Ltd|Company|Co)\b)'
                for match in re.finditer(consultant_pattern, text_content):
                    person_text = match.group(1)
                    org_text = match.group(2)
                    
                    person = get_or_create_entity(person_text, "PERSON")
                    org = get_or_create_entity(org_text, "ORGANIZATION")
                    
                    relation_type_raw = "consultant_for"
                    normalized_relation_type = self._normalize_relation_type(relation_type_raw)
                    relations.append({
                        "subject_id": person["id"],
                        "relation_type": normalized_relation_type,
                        "object_id": org["id"],
                        "properties": {
                            "confidence": 0.7,
                            "extraction_method": "regex",
                            "original_relation_text": relation_type_raw
                        }
                    })
            
            # Deduplicate relations
            unique_relations = {}
            for relation in relations:
                relation_key = f"{relation['subject_id']}_{relation['relation_type']}_{relation['object_id']}"
                if relation_key not in unique_relations or unique_relations[relation_key]['properties'].get('confidence', 0) < relation['properties'].get('confidence', 0):
                    unique_relations[relation_key] = relation
            
            # Return results
            if not processed_entities and not unique_relations:
                return {"success": True, "message": "No entities or relations found in the document.", "entities": [], "relations": []}

            return {"success": True, "entities": processed_entities, "relations": list(unique_relations.values())}
        except Exception as e:
            import traceback
            error_details = traceback.format_exc()
            return {
                "success": False, 
                "error": f"Error processing file for entity/relation extraction: {str(e)}", 
                "error_details": error_details,
                "entities": [], 
                "relations": []
            }

# Example usage (for testing, not part of the plugin itself)
if __name__ == "__main__":
    doc_processor = DocumentProcessor()
    # ... (rest of the example usage can remain for other functionalities)
    sample_file_path = "/home/ubuntu/sample_extraction_test.txt"
    with open(sample_file_path, "w") as f:
        f.write("Alice Smith works at Acme Corp. Bob Johnson is a consultant for Zeta Inc. Acme Corp is located in New York. Zeta Inc. is a competitor of Acme Corp.")
    
    print("\n--- Entity/Relation Extraction Test with Conceptual Disambiguation/Normalization ---")
    extraction_result = doc_processor.extract_entities_relations(sample_file_path)
    print(json.dumps(extraction_result, indent=2))
    if os.path.exists(sample_file_path):
        os.remove(sample_file_path)

